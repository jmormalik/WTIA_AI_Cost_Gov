#!/usr/bin/env python3
"""Build two English pitch decks (A: Pure Market Shift, B: Hero x Market Shift).

Charts are drawn with native python-pptx shapes (no matplotlib) so they stay
editable vectors and match the dark theme exactly. 4 visuals added 2026-07-31:
  - slide 9  Market      -> TAM/SAM/SOM concentric circles
  - slide 10 Competition -> 2x2 quadrant (after-the-fact<->pre-exec x visibility<->attribution)
  - villain slide        -> 41% non-LLM cost-composition bar
  - slide 16 Vision/Ask  -> revenue hockey-stick curve
Run:  python3 scripts/build_decks.py deliverables
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ----
BG      = RGBColor(0x0E, 0x14, 0x1B)   # near-black navy
PANEL   = RGBColor(0x16, 0x20, 0x2B)
ACCENT  = RGBColor(0xFF, 0x6B, 0x35)   # breaker/fire orange
ACCENT2 = RGBColor(0x35, 0xC9, 0xC9)   # teal
TEXT    = RGBColor(0xE8, 0xEC, 0xF1)
MUTED   = RGBColor(0x9A, 0xA7, 0xB4)
WARN    = RGBColor(0xF2, 0xC1, 0x4E)   # amber for hypotheses
GREY    = RGBColor(0x4A, 0x55, 0x60)   # LLM band (commodity grey)
RING1   = RGBColor(0x12, 0x20, 0x2B)   # TAM fill
RING2   = RGBColor(0x17, 0x32, 0x3C)   # SAM fill
QHILITE = RGBColor(0x14, 0x30, 0x3A)   # winning quadrant fill
AREA    = RGBColor(0x12, 0x3A, 0x40)   # hockey area fill

W, H = Inches(13.333), Inches(7.5)

def new_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs

def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s

def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf

def _set(p, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri"):
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    return r

def _accent_bar(slide, color=ACCENT):
    bar = slide.shapes.add_shape(1, Inches(0.7), Inches(0.7), Inches(0.16), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def _kicker(slide, text, num):
    tb, tf = _box(slide, Inches(0.95), Inches(0.62), Inches(9.5), Inches(0.5))
    p = tf.paragraphs[0]
    _set(p, text.upper(), 13, ACCENT2, bold=True)
    ntb, ntf = _box(slide, Inches(12.0), Inches(0.55), Inches(0.9), Inches(0.5))
    np = ntf.paragraphs[0]
    _set(np, num, 12, MUTED, bold=True, align=PP_ALIGN.RIGHT)

def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def _bullets(slide, items, top=Inches(2.35), left=Inches(0.95), width=Inches(11.4), size=17, gap=8):
    tb, tf = _box(slide, left, top, width, Inches(4.4))
    first = True
    for it in items:
        text = it[0]
        color = it[1] if len(it) > 1 and it[1] else TEXT
        bold = it[2] if len(it) > 2 else False
        indent = it[3] if len(it) > 3 else 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = indent
        bullet = "" if indent > 0 else "▸  "
        prefix = ("     " * indent) + ("–  " if indent > 0 else bullet)
        _set(p, prefix + text, size, color, bold=bold)
    return tb

# ---------- low-level chart primitives ----------
def _rect(slide, l, t, w, h, fill, line=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = line_w or Pt(1)
    sh.shadow.inherit = False
    return sh

def _oval(slide, cx, cy, d, fill=None, line=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - d/2), int(cy - d/2), int(d), int(d))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = line_w or Pt(1.25)
    sh.shadow.inherit = False
    return sh

def _line(slide, x1, y1, x2, y2, color=MUTED, w=Pt(1)):
    ln = slide.shapes.add_connector(2, int(x1), int(y1), int(x2), int(y2))
    ln.line.color.rgb = color; ln.line.width = w
    ln.shadow.inherit = False
    return ln

def _label(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=None):
    tb, tf = _box(slide, int(l), int(t), int(w), int(h))
    if anchor: tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    _set(tf.paragraphs[0], text, size, color, bold=bold, align=align)
    return tb

# ---------- the four charts ----------
def chart_tam_sam_som(slide):
    """Concentric circles. box right half."""
    cx = Inches(9.55); cy = Inches(4.55)
    _oval(slide, cx, cy, Inches(4.5), fill=RING1, line=ACCENT2, line_w=Pt(1.25))
    _oval(slide, cx, cy, Inches(2.95), fill=RING2, line=ACCENT2, line_w=Pt(1.0))
    _oval(slide, cx, cy, Inches(1.35), fill=ACCENT, line=None)
    # labels stacked down the vertical center
    _label(slide, cx-Inches(1.6), Inches(2.55), Inches(3.2), Inches(0.5),
           "TAM", 15, ACCENT2, bold=True, align=PP_ALIGN.CENTER)
    _label(slide, cx-Inches(1.6), Inches(2.90), Inches(3.2), Inches(0.4),
           "$0.5–1.8B  (2026)", 12, TEXT, align=PP_ALIGN.CENTER)
    _label(slide, cx-Inches(1.5), Inches(3.55), Inches(3.0), Inches(0.4),
           "SAM  $0.3–0.6B", 13, TEXT, bold=True, align=PP_ALIGN.CENTER)
    _label(slide, cx-Inches(0.85), Inches(4.28), Inches(1.7), Inches(0.6),
           "SOM", 13, BG, bold=True, align=PP_ALIGN.CENTER)
    _label(slide, cx-Inches(1.0), Inches(4.62), Inches(2.0), Inches(0.4),
           "$1.5–4.5M ARR", 10, BG, bold=True, align=PP_ALIGN.CENTER)
    _label(slide, Inches(6.95), Inches(6.75), Inches(5.6), Inches(0.4),
           "3-yr, not to scale  ·  main lens = governed AI spend × ~1–2% take-rate", 10, MUTED, align=PP_ALIGN.CENTER)

def chart_quadrant(slide):
    pl, pt = Inches(7.55), Inches(2.55)
    pw, ph = Inches(4.4), Inches(3.5)
    hx = int(pl + pw/2); vy = int(pt + ph/2)
    # winning quadrant highlight (top-right)
    _rect(slide, hx, int(pt), int(pw/2), int(ph/2), QHILITE)
    # axes
    _line(slide, pl, vy, int(pl+pw), vy, MUTED, Pt(1.25))
    _line(slide, hx, int(pt), hx, int(pt+ph), MUTED, Pt(1.25))
    # axis end labels
    _label(slide, int(pl-Inches(0.55)), vy-Inches(0.18), Inches(1.4), Inches(0.35), "After-the-fact", 9, MUTED)
    _label(slide, int(pl+pw-Inches(1.0)), vy-Inches(0.18), Inches(1.2), Inches(0.35), "Pre-execution", 9, MUTED, align=PP_ALIGN.RIGHT)
    _label(slide, hx-Inches(1.55), int(pt-Inches(0.34)), Inches(3.1), Inches(0.3), "Audit-grade attribution", 9, ACCENT2, bold=True, align=PP_ALIGN.CENTER)
    _label(slide, hx-Inches(1.4), int(pt+ph+Inches(0.03)), Inches(2.8), Inches(0.3), "Total visibility only", 9, MUTED, align=PP_ALIGN.CENTER)
    # players: (xfrac, yfrac from bottom-left, label, color, bold)
    players = [
        (0.18, 0.22, "Langfuse · Helicone", MUTED, False),
        (0.10, 0.06, "Logs · scripts", MUTED, False),
        (0.24, 0.74, "CloudZero · Vantage*", MUTED, False),
        (0.72, 0.16, "LiteLLM · Portkey", MUTED, False),
        (0.74, 0.80, "SatGate (pre-launch)", WARN, False),
        (0.88, 0.92, "Kyber", ACCENT, True),
    ]
    for xf, yf, lab, col, bold in players:
        px = int(pl + xf*pw); py = int(pt + (1-yf)*ph)
        d = Inches(0.20) if bold else Inches(0.14)
        _oval(slide, px, py, d, fill=col)
        # label placed to keep inside box
        lw = Inches(2.1)
        lx = px + Inches(0.14)
        if xf > 0.55: lx = px - lw - Inches(0.12)
        al = PP_ALIGN.RIGHT if xf > 0.55 else PP_ALIGN.LEFT
        _label(slide, int(lx), py-Inches(0.16), lw, Inches(0.34), lab, 10 if bold else 9,
               col if col != MUTED else TEXT, bold=bold, align=al)
    _label(slide, Inches(6.95), Inches(6.78), Inches(5.6), Inches(0.35),
           "*FinOps attributes cloud spend — off the agent path, after-the-fact", 9, MUTED, align=PP_ALIGN.CENTER)

def chart_cost_bar(slide):
    _label(slide, Inches(7.0), Inches(2.35), Inches(5.4), Inches(1.0), "41%", 60, ACCENT2, bold=True)
    _label(slide, Inches(7.05), Inches(3.35), Inches(5.4), Inches(0.4), "of COGS never touches a model.", 15, TEXT, bold=True)
    _label(slide, Inches(7.05), Inches(3.72), Inches(5.4), Inches(0.4), "It is tool & retrieval spend token dashboards don't see.", 11, MUTED)
    # segmented bar
    bl, bt, bw, bh = Inches(7.05), Inches(4.45), Inches(5.15), Inches(0.62)
    segs = [(0.59, GREY, "LLM API"), (0.24, ACCENT, "Search"), (0.17, ACCENT2, "External DB")]
    x = int(bl)
    for frac, col, _lab in segs:
        w = int(bw*frac)
        _rect(slide, x, int(bt), w, int(bh), col)
        x += w
    # legend
    leg = [("LLM API", "59% · $68.8K", GREY), ("Search Tool", "24% · $27.9K", ACCENT), ("External DB", "17% · $20.0K", ACCENT2)]
    ly = Inches(5.35)
    for i, (name, val, col) in enumerate(leg):
        y = int(ly + Inches(0.42)*i)
        _oval(slide, int(Inches(7.15)), int(y+Inches(0.12)), Inches(0.14), fill=col)
        _label(slide, Inches(7.4), y, Inches(2.3), Inches(0.35), name, 12, TEXT, bold=True)
        _label(slide, Inches(9.6), y, Inches(2.6), Inches(0.35), val, 12, MUTED)
    _label(slide, Inches(7.05), Inches(6.75), Inches(5.4), Inches(0.4),
           "Search + DB = $47.8K nobody attributes to the agent that spent it.", 10, ACCENT2, align=PP_ALIGN.LEFT)

def chart_hockey(slide):
    pl, pr = int(Inches(7.55)), int(Inches(12.25))
    y_base, y_top = int(Inches(6.05)), int(Inches(2.75))
    pw = pr - pl; ph = y_base - y_top
    vals = [0.0, 0.24, 1.35, 4.8]; vmax = 4.8
    xs = [int(pl + i/3*pw) for i in range(4)]
    ys = [int(y_base - (v/vmax)*ph) for v in vals]
    # axes
    _line(slide, pl, y_base, pr, y_base, MUTED, Pt(1))
    _line(slide, pl, y_top, pl, y_base, MUTED, Pt(1))
    # area under curve (freeform)
    fb = slide.shapes.build_freeform(xs[0], y_base, scale=1)
    pts = [(xs[i], ys[i]) for i in range(4)] + [(xs[3], y_base)]
    fb.add_line_segments(pts, close=True)
    area = fb.convert_to_shape()
    area.fill.solid(); area.fill.fore_color.rgb = AREA
    area.line.color.rgb = ACCENT2; area.line.width = Pt(2.5)
    area.shadow.inherit = False
    # markers + value labels + customer counts
    val_txt = ["", "$0.24M", "$1.35M", "$4.8M"]
    cust = ["", "8", "45", "160"]
    yr = ["Yr0", "Yr1", "Yr2", "Yr3"]
    for i in range(4):
        if i > 0:
            _oval(slide, xs[i], ys[i], Inches(0.16), fill=ACCENT)
            _label(slide, xs[i]-Inches(0.85), ys[i]-Inches(0.46), Inches(1.7), Inches(0.35),
                   val_txt[i], 13, TEXT, bold=True, align=PP_ALIGN.CENTER)
        _label(slide, xs[i]-Inches(0.6), y_base+Inches(0.06), Inches(1.2), Inches(0.3),
               yr[i], 11, MUTED, align=PP_ALIGN.CENTER)
        if cust[i]:
            _label(slide, xs[i]-Inches(0.7), y_base+Inches(0.36), Inches(1.4), Inches(0.3),
                   cust[i]+" cust.", 10, ACCENT2, align=PP_ALIGN.CENTER)
    _label(slide, Inches(6.95), Inches(6.95), Inches(5.6), Inches(0.4),
           "ARR @ $30K ACV  ·  ⚠ strong-execution case (target NRR ~110%)", 10, WARN, align=PP_ALIGN.CENTER)

# ---------- slide builders ----------
def title_slide(prs, product, tagline, sub, team, version_label):
    s = _blank(prs)
    blk = s.shapes.add_shape(1, Inches(0.0), Inches(2.55), Inches(0.35), Inches(2.4))
    blk.fill.solid(); blk.fill.fore_color.rgb = ACCENT; blk.line.fill.background()
    tb, tf = _box(s, Inches(0.95), Inches(2.4), Inches(11.4), Inches(3.0))
    p = tf.paragraphs[0]
    _set(p, product, 40, TEXT, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(6)
    _set(p2, "Make the cost of every AI agent explainable.", 22, ACCENT2, bold=False)
    p3 = tf.add_paragraph(); p3.space_before = Pt(22)
    _set(p3, tagline, 20, TEXT, italic=True)
    p4 = tf.add_paragraph(); p4.space_before = Pt(10)
    _set(p4, sub, 14, MUTED)
    ftb, ftf = _box(s, Inches(0.95), Inches(6.5), Inches(11.4), Inches(0.6))
    fp = ftf.paragraphs[0]
    _set(fp, team, 13, MUTED)
    fp2 = ftf.add_paragraph()
    _set(fp2, version_label, 11, ACCENT, bold=True)

def content_slide(prs, kicker, num, headline, bullets, note, headline_color=TEXT, hsize=27):
    s = _blank(prs)
    _accent_bar(s)
    _kicker(s, kicker, num)
    tb, tf = _box(s, Inches(0.95), Inches(1.15), Inches(11.6), Inches(1.4))
    p = tf.paragraphs[0]
    _set(p, headline, hsize, headline_color, bold=True)
    _bullets(s, bullets)
    _notes(s, note)
    return s

def content_slide_chart(prs, kicker, num, headline, bullets, chart_fn, note, headline_color=TEXT, hsize=25):
    """Two-column: trimmed bullets left (~5.7in), chart drawn on the right."""
    s = _blank(prs)
    _accent_bar(s)
    _kicker(s, kicker, num)
    tb, tf = _box(s, Inches(0.95), Inches(1.15), Inches(11.6), Inches(1.4))
    _set(tf.paragraphs[0], headline, hsize, headline_color, bold=True)
    _bullets(s, bullets, top=Inches(2.35), left=Inches(0.95), width=Inches(5.75), size=15, gap=9)
    chart_fn(s)
    _notes(s, note)
    return s

# =====================================================================
# SHARED BODY (slides 6-13, 15)
# =====================================================================
def add_shared_body(prs):
    # [6] Vehicle
    content_slide(prs, "Vehicle · The Solution", "6",
        "A layer that attributes and controls cost — on the execution path.",
        [
            ("Logging gives you the total; we give you the breakdown. We attribute & govern every LLM + tool + external-API call at request time — per feature / team / customer.",),
            ("\"Not CCTV after the fire — a breaker before it ignites.\" Block before execution + audit-grade reason for the block.", ACCENT2, True),
            ("Enter via a Trust Ladder:", TEXT, True),
            ("Step 1 — off-path reports (no proxy)", MUTED, False, 1),
            ("Step 2 — async policy", MUTED, False, 1),
            ("Step 3 — block inside the customer-owned gateway (we never own the path)", MUTED, False, 1),
        ],
        "Trust Ladder = our risk answer. Directly reflects the 'only OSS on the critical path' signal. Team-confirmed 7/23.")

    # [7] How it works / Demo
    content_slide(prs, "How It Works · Demo", "7",
        "Three mechanisms — all in service of real-time cost attribution & control.",
        [
            ("1. Inline policy engine — evaluate & block per session / customer / feature budget BEFORE execution. Block reason to an audit log. (SPOF/latency handled via sidecar, async, bypass.)",),
            ("2. Semantic usage attribution — tag every call request → agent step → LLM / tool / external API. Keep a non-LLM unit-price catalog to sum true total cost.",),
            ("3. Dynamic model routing (by-product) — route by complexity, save with no quality loss. Mentioned only as a bonus (public technique = not a moat).", MUTED, False),
        ],
        "Demo slide = embed the React demo (per-customer unit economics + real-time budget hard-limits + routing). Never use the 'kernel/C++' framing from the old Gemini draft.")

    # [8] Customer
    content_slide(prs, "Customer · Who Buys", "8",
        "Can't build like Meta — but far more spend & audit pressure than a startup.",
        [
            ("Paid target: mid-market / scale-up B2B — $50K–$500K/mo AI spend, 50–300 people, Series A–C. Not big enough to staff an infra team → can't build, must buy.", TEXT, True),
            ("Why B2B: per-customer unit cost is a B2B need (customer = contract). B2C rolls cost up by category — a B2B CTO must defend margin per account. (An ML lead we spoke with framed it exactly this way.)", ACCENT2, True),
            ("Tooling / spend = 1–3% (clear ROI). Skips the death valley (small = negative ROI / enterprise = builds its own). Buyer = CTO, champion = AI lead, finance = influencer.",),
            ("Korean AI startups = free design partners (data & references); paid = US mid-market.", MUTED, False),
            ("⚠  Bet, not fact: 'can't build → must buy' + the spend threshold & buyer are what we are validating in interviews.", WARN, False),
        ],
        "This slide pre-empts two rebuttals ('just buy tokens' / 'build it yourself'). Label the 'can't build' load-bearing claim as a ⚠ falsifiable bet, and pre-empt the Liner counter — 'estimable ≠ audit-grade attribution'. B2B/per-customer rationale added 7/29.")

    # [9] Market  -- TAM/SAM/SOM concentric circles
    content_slide_chart(prs, "Market · TAM / SAM / SOM", "9",
        "Korea is the wedge; the market is global agent infra. Bet the slope, not the size.",
        [
            ("Lens A (governed AI spend × ~1–2% take-rate) = main; Lens B (AI FinOps + LLMOps) = cross-check. CAGR ~20%.",),
            ("ACV $30K = median of competitor mid-market ACVs.",),
            ("SOM engine: 50–150 US paid, seeded by the Korea wedge + US network.", ACCENT2, True),
            ("⚠  SAM population (10–20K firms) = hypothesis, refined by interview screening pass-rate.", WARN, False),
        ],
        chart_tam_sam_som,
        "Never present the old LAM $80K as market size. Take-rate anchor CORRECTED to ~1-2% (deep research 7/28). We don't price on take-rate (per-agent) — it's a market-sizing sanity check only. Source: Market/TAM-SAM-SOM.md.")

    # [10] Competition  -- 2x2 quadrant
    content_slide_chart(prs, "Competition", "10",
        "Two axes: after-the-fact vs. pre-execution × total visibility vs. audit-grade attribution.",
        [
            ("Showing the total is table stakes — logging, observability & billing all do it.", TEXT, True),
            ("The gap = attributing the agent's full cost per feature / team / customer, audit-grade, before execution.", ACCENT2, True),
            ("FinOps (CloudZero · Vantage) does per-unit attribution — but for cloud spend, after-the-fact, off the agent path.",),
            ("Our square — pre-execution × audit-grade attribution. Nobody occupies it.", TEXT, True),
        ],
        chart_quadrant,
        "Reframed 7/29: don't say rivals 'can't see cost' — they show totals. HONESTY: FinOps already sells per-unit attribution, but for cloud/infra spend, after-the-fact, not the agent execution path. Our differentiation = the intersection. SatGate pre-launch = timing validation, not absent demand.")

    # [11] Moat
    content_slide(prs, "Moat", "11",
        "What gets copied is features. What can't be caught up is data.",
        [
            ("Caching, routing, hard-limits = replicable in 6 months. Not a moat.", MUTED, False),
            ("Moat = things that get stronger as customer data accumulates:", TEXT, True),
            ("1. Unit-economics accuracy gap — more traffic → sharper runaway detection & cost simulation. Followers start from zero data.", None, False, 1),
            ("2. Peer benchmark (lagging) — 'top 20% cost-per-request vs. your peer type.' Scales with customer count; structurally impossible for latecomers.", None, False, 1),
            ("Once this cost view is embedded in pricing & board routines, switching cost spikes.", ACCENT2, True),
        ],
        "Timeline: entry (free credit-countdown wedge) → differentiation (unit economics + unit-price catalog) → moat (simulation accuracy + benchmark).")

    # [12] BM / Pricing
    content_slide(prs, "Business Model · Pricing", "12",
        "A company that sells cost does not earn in proportion to the customer's cost.",
        [
            ("Spend-proportional pricing ✗ (conflicts with the savings incentive). Price on per-agent (# of monitored agents) → incentives aligned.", TEXT, True),
            ("3 tiers: Free (monitoring/blocking = commodity) → Growth ~$18K/yr (automated COGS report = paid entry) → Enterprise ~$60K/yr (audit-grade + in-path blocking). Blended ACV ~$30K.",),
            ("Anchored in comps: LLM observability $5–25K (Langfuse/Helicone/Portkey) < us $30K < cloud FinOps $30–130K (CloudZero/Cloudability). Precedent for per-unit pricing: Agentforce, Datadog per-host.", MUTED, False),
            ("Expansion = agents 20 → 40 → 80 grows revenue = net-negative churn.", ACCENT2, True),
            ("Slogan: \"Sell the ledger before you sell the breaker.\"", None, True),
            ("⚠  Price figures are hypotheses to validate in interviews.", WARN, False),
        ],
        "'Flat or usage?' → 'Subscription-led, but the usage axis is agent count, not spend.' ACV $30K = median of competitor mid-market ACVs (deep research 7/28). Detail: Solution/가격-재무-딥리서치-2026-07.md.")

    # [13] GTM
    content_slide(prs, "Go-To-Market", "13",
        "The champion bites for free; the CTO opens the wallet for COGS.",
        [
            ("Landing: free Cost Autopsy report (off-path) → AI lead adopts bottom-up.",),
            ("Conversion: at an incident or credit burn-out, the champion pitches the CTO paid → land with COGS & margin reports.",),
            ("Market order: Korea (free design partners / validation) → US mid-market (paid).", ACCENT2, True),
            ("Blocker pre-empt: to IT/legal — 'metadata only + cross-vendor independent audit.'",),
        ],
        "Trust Ladder steps 1→2→3 ARE the GTM expansion path. Based on the stakeholder-map action plan.")

    # [14] Traction
    content_slide(prs, "Traction", "14",
        "In place of traction, how precisely we've narrowed the problem — and our first real signal.",
        [
            ("How we sell = we don't sell to everyone. A discipline that quickly filters anti-ICP.", TEXT, True),
            ("First real signal — an ML lead at a Korean AI-search company (2 rounds) confirmed build/buy splits cost structure, and that per-customer attribution needs deliberate data work.", ACCENT2, False),
            ("Outreach sent (Korea AI startups + US mid-market Top 5). Design partners: N in progress.",),
            ("Formal interviews: 0 → N in progress, first signal secured (pricing built on data, honestly stated).",),
            ("Target metric = net-negative churn (agent expansion).", None, True),
        ],
        "Don't hide the 0 formal interviews — frame as 'building pricing on data.' Fill with real interview/design-partner counts before 8/6.")

    # [15] Team
    content_slide(prs, "Team", "15",
        "Technology × customer development — a 2-person core.",
        [
            ("Taewon — technology & product (agent orchestration · infrastructure).",),
            ("Jaemin (Malik) — CEO · customer development · GTM · research.",),
            ("(Advisors & design-partner relationships added as secured.)", MUTED, False),
            ("⚠  Confirm role split — deck vs. memory mismatch on who is CEO; align before the pitch.", WARN, False),
        ],
        "Keep it short. One line on 'why us' to be filled by the team. Resolve the Team-slide role inconsistency before the pitch.")

# =====================================================================
# VERSION A — Pure Market Shift
# =====================================================================
def build_version_a(path):
    prs = new_deck()
    title_slide(prs, "Kyber",
        "\"Not CCTV after the fire — a breaker before it ignites.\"",
        "Govern the LLM / tool / external-API cost agents burn on their own — on the execution path — and explain it, audit-grade.",
        "Team · 2026 · Presenter",
        "VERSION A — Pure Market Shift")

    content_slide(prs, "The Shift · The spine (undeniable layer)", "2",
        "The decider of AI spend has moved from humans to the model.",
        [
            ("Old world: cost is declared by the developer in code → predictable, controllable.",),
            ("New world: the model itself decides which tool to call, when, and how many times → the spend trigger moves to the model.", TEXT, True),
            ("Production-agent tokens ≈ 50× a chatbot. Finance has lost control.",),
            ("Why Now: someone has to solve this now, with or without us — Meta & Uber already built their own gateways; SatGate pre-launched.", ACCENT2, False),
        ],
        "The ONLY claim we make as 'the shift' is 'agents spend on their own' — observable, undeniable. Everything contentious is downgraded to ⚠ bets later.")

    # [3] Villain -- 41% non-LLM cost bar
    content_slide_chart(prs, "The Villain · Loser's tools", "3",
        "The total is visible. The breakdown isn't.",
        [
            ("Cloud logging already shows the monthly total — that's not the gap. No one can break it down by feature / team / customer.", TEXT, True),
            ("Observability = monthly, after-the-fact. Static keys & caps = a blunt knife that can't tell normal from runaway.",),
            ("And tools see only up to the LLM call. Search, crawling, third-party APIs are off-screen — so true unit cost is never attributed.", ACCENT2, True),
        ],
        chart_cost_bar,
        "Reframed 7/29: DON'T claim 'cost is invisible' — cloud logging shows the total. The real gap = per-feature/team/customer attribution + the 41% non-LLM spend nobody attributes. Numbers match the React demo for cross-consistency.")

    content_slide(prs, "Real Pain · Evidence, not a hero", "4",
        "The problem isn't 'we don't know how much' — it's 'we can't explain why.'",
        [
            ("Acute (verified wedge): one dev, an infinite retry loop → $47K billed, unnoticed for 11 days. A B2B SaaS: API usage +400% overnight.", TEXT, True),
            ("At enterprise scale, control failure = hundreds of millions/yr.", MUTED, False),
            ("Redefinition: the spike is just the trigger. The real loss = even after the incident, no answer to 'cost per customer / which feature eats the margin.'",),
            ("⚠  Chronic (always-on COGS explanation) = an expansion hypothesis we are validating in interviews.", WARN, False),
        ],
        "Placed as EVIDENCE of the 'loser,' not a personal hero story. Split acute (verified) vs. chronic (⚠ expansion bet). Source: why-now.md.")

    content_slide(prs, "The Promised Land · The winner", "5",
        "You become a company that can explain AI cost per customer & per feature — in real time.",
        [
            ("'We cut your AI bill' ✗ → 'We let you explain your AI unit cost' ✓.", ACCENT2, True),
            ("Pricing, the board, and due diligence become real-time numbers instead of guesses.",),
        ],
        "Here we explicitly drop 'savings.' Savings appears later (slide 7) only as a by-product.")

    add_shared_body(prs)

    # [16] Vision / Ask -- hockey stick
    content_slide_chart(prs, "Vision · Ask", "16",
        "As agents grow, revenue bends into a curve.",
        [
            ("Revenue only (no cost/net income). Engine = new logos + agent expansion.", MUTED, False),
            ("Vision: in an age where autonomous software spends money, a world where every dollar is explainable.", ACCENT2, True),
            ("Our ask isn't a check — it's connections:", ACCENT, True),
            ("(a) AI / infra advisors to pressure-test the wedge  ·  (b) mid-market teams to onboard as design partners  ·  (c) warm intros to production-agent CTOs & AI leads.", ACCENT2, False),
        ],
        chart_hockey,
        "Finance goal = a 2nd meeting, not a precise forecast. THE ASK IS INTRODUCTIONS, NOT CAPITAL. Curve is credible but AGGRESSIVE (5.6x then 3.6x) — strong-execution case; NRR 110% = top-quartile target; ACV $30K = competitor-median validated. Don't raise a funding amount unless a judge asks.")

    prs.save(path)
    print("saved", path, "slides:", len(prs.slides._sldIdLst))

# =====================================================================
# VERSION B — Hero Journey × Market Shift
# =====================================================================
def build_version_b(path):
    prs = new_deck()
    title_slide(prs, "Kyber",
        "\"Not CCTV after the fire — a breaker before it ignites.\"",
        "Govern the LLM / tool / external-API cost agents burn on their own — on the execution path — and explain it, audit-grade.",
        "Team · 2026 · Presenter",
        "VERSION B — Hero Journey × Market Shift")

    content_slide(prs, "Cold Open · Meet Sarah (blink test)", "2",
        "Meet Sarah. CTO of a mid-market SaaS.",
        [
            ("Sarah's agents take live customer traffic in production. They call search, tools, and external APIs on their own.",),
            ("Friday night, one retry loop quietly burned $47K over 11 days. No dashboard caught it in real time.", TEXT, True),
            ("Monday board meeting: 'What's our AI cost per customer?' — Sarah couldn't answer.", ACCENT2, True),
            ("⚠  Sarah = an illustrative representative persona, not a real interviewed customer.", WARN, False),
        ],
        "Opening with a specific person conveys 'this is for a specific person (a mid-market CTO), not everyone.' The 'representative persona' label is mandatory.")

    content_slide(prs, "Zoom Out · This isn't Sarah alone", "3",
        "Not one person's mistake — the decider of spend has moved from humans to the model.",
        [
            ("Old world: developers declare cost in code → predictable.",),
            ("New world: the model itself pulls the spend trigger. Production tokens ≈ 50× a chatbot.", TEXT, True),
            ("Companies that adapt explain their cost; companies that don't lose their margin.",),
            ("Why Now: Meta & Uber built their own gateways, SatGate pre-launched — this must be solved now, with or without us.", ACCENT2, False),
        ],
        "Bridge slide hero → market shift. Only here assert the shift as 'one undeniable layer.'")

    # [4] Villain -- 41% non-LLM cost bar
    content_slide_chart(prs, "The Villain · Why Sarah couldn't answer", "4",
        "Sarah's dashboard shows the total. It can't show the breakdown.",
        [
            ("Cloud logging gave Sarah the monthly total — but not 'per customer, per feature.' That split needs deliberate data structuring no tool did for her.", TEXT, True),
            ("Observability = after-the-fact. Static keys & caps = a blunt knife.",),
            ("Tools see only up to the LLM call — search, crawling, third-party APIs off-screen. So Sarah saw how much, but never to whom.", ACCENT2, True),
        ],
        chart_cost_bar,
        "Reframed 7/29: the gap isn't 'total invisible' (cloud logging shows it) — it's per-feature/team/customer attribution + the 41% non-LLM spend. Keep incident numbers ($47K) tied to Sarah.")

    content_slide(prs, "The Promised Land · Sarah transformed", "5",
        "Now Sarah answers the Monday board question in 10 seconds.",
        [
            ("AI cost per customer & per feature, as real-time numbers.",),
            ("'We cut your bill' ✗ → 'We let you explain your unit cost' ✓.", ACCENT2, True),
            ("Pricing, the board, and due diligence become measured, not guessed.",),
        ],
        "Completes the hero journey (problem → experiencing the solution → change).")

    add_shared_body(prs)

    # [16] Vision / Ask -- hockey stick
    content_slide_chart(prs, "Vision · Ask", "16",
        "As agents grow, revenue bends into a curve.",
        [
            ("Revenue only (no cost/net income). Engine = new logos + agent expansion.", MUTED, False),
            ("Vision: in an age where autonomous software spends money, a world where every dollar is explainable.", ACCENT2, True),
            ("Our ask isn't a check — it's connections:", ACCENT, True),
            ("(a) AI / infra advisors  ·  (b) mid-market design partners  ·  (c) warm intros to CTOs like Sarah running agents in production.", ACCENT2, False),
        ],
        chart_hockey,
        "Finance goal = a 2nd meeting. THE ASK IS INTRODUCTIONS, NOT CAPITAL. Curve credible but AGGRESSIVE; NRR 110% = top-quartile target; ACV $30K = competitor-median validated. Version B recovers the Sarah persona (her world grows). Don't name a funding amount unless a judge asks.")

    prs.save(path)
    print("saved", path, "slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    import sys
    out = sys.argv[1] if len(sys.argv) > 1 else "."
    build_version_a(f"{out}/Pitch_Deck_VersionA_Market-Shift_EN.pptx")
    build_version_b(f"{out}/Pitch_Deck_VersionB_Hero-Hybrid_EN.pptx")
